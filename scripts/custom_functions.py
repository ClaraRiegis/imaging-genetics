#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Created on Thu May 30 13:05:33 2024

@author: clarariegis
"""
import os
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import numpy as np
import matplotlib.pyplot as plt
from collections import defaultdict
from functools import reduce


#%%





def create_subplot(num_plots, max_col = 4):
    # Calculate the number of rows and columns based on the number of plots
    num_cols = min(num_plots, max_col)  # Maximum of 5 plots per row
    num_rows = (num_plots + num_cols - 1) // num_cols  # Ceiling division to ensure at least one row
    
    # Calculate the figure size to maintain the desired aspect ratio
    aspect_ratio = 55/65  # Desired height-to-width ratio
    width = 6  # Width of each subplot in inches
    height = width * aspect_ratio * num_rows / num_cols  # Calculate height based on aspect ratio
    
    # Create the subplot
    fig, axs = plt.subplots(num_rows, num_cols, figsize=(width, height), dpi=1500)
    
    # Flatten the axs array if it's more than one-dimensional
    if num_rows == 1:
        axs = axs[np.newaxis, :]
    axs = axs.flatten()

    
    return fig, axs



def create_directory_if_not_exists(directory_path):
    """
    Create a directory if it does not exist.
    
    Parameters:
    directory_path (str): Path to the directory to be created.
    """
    if not os.path.exists(directory_path):
        os.makedirs(directory_path, exist_ok=True)
        print(f"Directory created: {directory_path}")
    else:
        print(f"Directory already exists: {directory_path}")


def format_number(value, num_format):
    """
    Format a number based on the specified format.
    
    Parameters:
    value: The number to format.
    num_format (str or int): The format to apply. If 'int', convert to integer. If 'float', leave as is. If an integer, round to that number of decimal places.
    
    Returns:
    str: The formatted number as a string.
    """
    if num_format == 'int':
        return str(int(value))
    elif num_format == 'float':
        return str(value)
    elif isinstance(num_format, int):
        return str(round(value, num_format))
    else:
        raise ValueError("num_format should be 'int', 'float', or an integer.")

def add_figures_to_presentation(input_ppt_path, scale_ppt_path, fig_paths, desired_widths, desired_heights, left_positions, top_positions, output_ppt_path, num_format='float'):
    """
    Add figures and tables to a PowerPoint presentation with specified positions and dimensions.

    This function loads an existing PowerPoint presentation, adds slides with images or tables at specified
    positions and dimensions, and saves the modified presentation to a new file.

    Parameters:
    input_ppt_path (str): Path to the input PowerPoint file.
    scale_ppt_path (str): Path to the scale PowerPoint file (used for scaling the images).
    fig_paths (list of list of str): A list of lists, where each sublist contains paths to images or tables for each slide.
    desired_widths (list of Inches): A list of desired widths for the images or tables.
    desired_heights (list of Inches): A list of desired heights for the images or tables.
    left_positions (list of Inches): A list of left positions for the images or tables on the slides.
    top_positions (list of Inches): A list of top positions for the images or tables on the slides.
    output_ppt_path (str): Path to save the modified PowerPoint file.
    num_format (str or int): Specifies the format for numbers in tables. 'int' for integers, 'float' to leave as is, or an integer to round to that many decimal places.

    Example:
    input_ppt_path = 'Figures/Slides/mphil_figures_clean.pptx'
    scale_ppt_path = 'Figures/Slides/scale.pptx'
    fig_paths = [['path/to/fig1_1.png', 'path/to/fig1_2.png'], ['path/to/fig2_1.csv', 'path/to/fig2_2.csv'], ['path/to/fig3_1.png', 'path/to/fig3_2.png']]
    desired_widths = [Inches(7), Inches(3.5), Inches(3.5)]
    desired_heights = [Inches(5), Inches(2.5), Inches(2.5)]
    left_positions = [Inches(1), Inches(0.75), Inches(3.25)]
    top_positions = [Inches(0.5), Inches(5), Inches(5)]
    output_ppt_path = 'Figures/Slides/mphil_figures_clean2.pptx'
    num_format = 2

    add_figures_to_presentation(input_ppt_path, scale_ppt_path, fig_paths, desired_widths, desired_heights, left_positions, top_positions, output_ppt_path, num_format)
    """
    # Load existing PowerPoint presentations
    prs = Presentation(input_ppt_path)
    
    # Check if the scale presentation exists; if not, create it
    if not os.path.exists(scale_ppt_path):
        scale_ppt_dir = os.path.dirname(output_ppt_path)
        scale_ppt_path = os.path.join(scale_ppt_dir, os.path.basename(scale_ppt_path))
        sc_prs = Presentation()
        sc_prs.save(scale_ppt_path)
        print(f"Created new scale presentation at {scale_ppt_path}")
    else:
        sc_prs = Presentation(scale_ppt_path)

    # Loop through each figure path set
    for i in range(len(fig_paths[0])):
        slide_layout = prs.slide_layouts[6]  # Use layout with title and content
        slide = prs.slides.add_slide(slide_layout)
        sc_slide = sc_prs.slides.add_slide(slide_layout)
        
        # Add images or tables to the slide
        for h in range(len(fig_paths)):
            file_path = fig_paths[h][i]
            if file_path.endswith('.csv'):
                # Add table from CSV
                df = pd.read_csv(file_path)
                rows, cols = df.shape
                left, top = left_positions[h], top_positions[h]
                width, height = desired_widths[h], desired_heights[h]

                table = slide.shapes.add_table(rows+1, cols, left, top, width, height).table

                # Set column widths
                col_width = int(width / cols)
                for col in range(cols):
                    table.columns[col].width = col_width

                # Write column headings
                for col, heading in enumerate(df.columns):
                    cell = table.cell(0, col)
                    cell.text = heading
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.bold = True
                        paragraph.font.size = Pt(12)
                        paragraph.alignment = PP_ALIGN.CENTER

                # Write data
                for row in range(rows):
                    for col in range(cols):
                        cell = table.cell(row+1, col)
                        cell.text = format_number(df.iloc[row, col], num_format)
                        for paragraph in cell.text_frame.paragraphs:
                            paragraph.font.size = Pt(10)
                            paragraph.alignment = PP_ALIGN.CENTER
            else:
                # Get the original width and height of the image
                image = sc_slide.shapes.add_picture(file_path, Inches(0), Inches(0))
                original_width, original_height = image.image.size
                
                # Calculate the scaling factor
                scaling_factor = min(desired_widths[h] / original_width, desired_heights[h] / original_height)
                
                # Calculate the new width and height
                new_width = int(original_width * scaling_factor)
                new_height = int(original_height * scaling_factor)
                
                # Add the image to the slide
                slide.shapes.add_picture(file_path, left_positions[h], top_positions[h], width=new_width, height=new_height)
    
    # Save the modified presentation
    prs.save(output_ppt_path)
    print(f"Presentation saved as {output_ppt_path}")
    
    

def save_cross_sectional_measures(df, quant_covar, cat_covar, roi_names, dir_path, meas):
    """
    Save cross-sectional measures for specified time points.

    Parameters:
    df (pd.DataFrame): DataFrame containing the data.
    quant_covar (pd.DataFrame): DataFrame containing quantitative covariates.
    cat_covar (pd.DataFrame): DataFrame containing categorical covariates.
    roi_names (list): List of ROI names.
    dir_path (str): Directory path to save the output files.
    meas (str): Measurement name to be included in the file name.

    Returns:
    None
    """
    # Time points to compare
    time_points = ['t1', 't2', 't3']
    
    variables = ['fam_id', 'ID'] + roi_names

    for time_point in time_points:
        # Save cross-sectional measures
        cross_sect_dfs = df.loc[df['eventname'] == time_point].dropna(subset=variables)
        cross_sect_dfs[variables].to_csv(f"{dir_path}/meas_{meas}_{time_point}.txt", index=False, sep=' ', header=False)
        
        # Filter quantitative and categorical covariates
        covar_quant = quant_covar[quant_covar['IID'].isin(cross_sect_dfs['ID'])]
        cat_covar = cat_covar[quant_covar['IID'].isin(cross_sect_dfs['ID'])]

        # Save covariates
        covar_quant.to_csv(f"{dir_path}/quant_{time_point}.txt", index=False, sep=' ', header=False)
        cat_covar.to_csv(f"{dir_path}/cat_{time_point}.txt", index=False, sep=' ', header=False)

# Example usage:
# df = pd.read_csv('path_to_your_dataframe.csv')
# quant_covar = pd.read_csv('path_to_quant_covar.csv')
# cat_covar = pd.read_csv('path_to_cat_covar.csv')
# roi_names = ['roi1', 'roi2', 'roi3']  # Example ROI names
# dir_path = 'your_directory_path'
# meas = 'your_measurement_name'

# save_cross_sectional_measures(df, quant_covar, cat_covar, roi_names, dir_path, meas)



import pandas as pd
from collections import defaultdict
from functools import reduce

def calculate_and_save_rate_of_change(df, roi_names, time_points, dir_path, meas, ids, quant_covar, cat_covar):
    """
    Calculate the rate of change for each ROI between specified time points and save the results to text files.
    Also saves quant_covar and cat_covar based on the IDs present in the rate of changes.

    Parameters:
    df (pd.DataFrame): DataFrame containing the data.
    roi_names (list): List of ROI names.
    time_points (list): List of time points to compare.
    dir_path (str): Directory path to save the output files.
    meas (str): Measurement name to be included in the file name.
    ids (pd.DataFrame): DataFrame containing the ID column for merging.
    quant_covar (pd.DataFrame): DataFrame containing quantitative covariates.
    cat_covar (pd.DataFrame): DataFrame containing categorical covariates.

    Returns:
    None
    """
    # Pivot the dataframe to get separate columns for each time point
    df_pivot = df.pivot(index='ID', columns='eventname', values=roi_names + ['age_scaled'])

    # Flatten the MultiIndex columns
    df_pivot.columns = ['_'.join(col).strip() for col in df_pivot.columns.values]

    # Initialize a dictionary to store rate of change dataframes
    rate_of_change_dfs = defaultdict(list)

    # Loop through each value column
    for value_col in roi_names:
        # Loop through pairs of time points to calculate rate of change
        for i in range(len(time_points) - 1):
            for j in range(i + 1, len(time_points)):
                t1, t2 = time_points[i], time_points[j]
                rate_column_name = f'rate_{value_col}_{t1}{t2}'
                rate_df_name = f'rate_{t1}{t2}'

                # Calculate the rate of change
                df_pivot[rate_column_name] = (df_pivot[f'{value_col}_{t2}'] - df_pivot[f'{value_col}_{t1}']) / (df_pivot[f'age_scaled_{t2}'] - df_pivot[f'age_scaled_{t1}'])

                # Drop rows where necessary time points are missing
                rate_df = df_pivot.dropna(subset=[f'{value_col}_{t1}', f'{value_col}_{t2}', f'age_scaled_{t1}', f'age_scaled_{t2}'])

                # Select relevant columns
                rate_df = rate_df[[rate_column_name]]

                # Store the result dataframe in the dictionary
                rate_of_change_dfs[rate_df_name].append(rate_df)

    # Merge and save the dataframes for each pair of time points
    for time_pair in ['t1t2', 't1t3', 't2t3']:
        df_merged = reduce(lambda left, right: pd.merge(left, right, on='ID', how='outer'), rate_of_change_dfs[f'rate_{time_pair}']).dropna()
        df_merged = pd.merge(ids, df_merged, on="ID")
        
        # Save the rate of change dataframe
        df_merged.to_csv(f"{dir_path}/meas_{meas}_{time_pair}.txt", sep=' ', index=False, header=False)
        
        # Save the quant_covar and cat_covar based on the IDs present in the rate of changes
        covar_quant = quant_covar[quant_covar['IID'].isin(df_merged['ID'])]
        # Drop columns with 'age' in their names
        covar_quant = covar_quant.loc[:, ~covar_quant.columns.str.contains('age', case=False)]

        cat_covar = cat_covar[cat_covar['IID'].isin(df_merged['ID'])]
        
        covar_quant.to_csv(f"{dir_path}/quant_{time_pair}.txt", sep=' ', index=False, header=False)
        cat_covar.to_csv(f"{dir_path}/cat_{time_pair}.txt", sep=' ', index=False, header=False)

# Example usage:
# df = pd.read_csv('path_to_your_dataframe.csv')
# roi_names = ['roi1', 'roi2', 'roi3']  # Example ROI names
# time_points = ['t1', 't2', 't3']
# dir_path = 'your_directory_path'
# meas = 'your_measurement_name'
# ids = pd.read_csv('path_to_ids.csv')  # Example IDs DataFrame
# quant_covar = pd.read_csv('path_to_quant_covar.csv')  # Example quantitative covariates DataFrame
# cat_covar = pd.read_csv('path_to_cat_covar.csv')  # Example categorical covariates DataFrame

# calculate_and_save_rate_of_change(df, roi_names, time_points, dir_path, meas, ids, quant_covar, cat_covar)




def is_convertible_to_numeric(series):
    try:
        pd.to_numeric(series)
        return True
    except ValueError:
        return False

# # Convert columns to numeric if possible
# for col in df.columns:
#     if is_convertible_to_numeric(df[col]):
#         df[col] = pd.to_numeric(df[col])






def get_lim(df, lim = None):
    """
    Determines the y-axis limit based on the minimum and maximum numeric values in a DataFrame.
    
    Parameters:
    df (pd.DataFrame): The DataFrame containing the measures.
    
    Returns:
    float: The y-axis limit.
    """
    # Select only numeric columns
    numeric_df = df.select_dtypes(include=[float, int])
    
    # Compute the minimum and maximum values from the numeric columns
    min_value = numeric_df.min().min()
    max_value = numeric_df.max().max()
    
    
    # Determine the y-axis limit
    if abs(min_value) > max_value:
        y_lim = abs(min_value)
    else:
        y_lim = max_value
    
    if lim == None: return y_lim
    elif lim == 'min': return min_value
    elif lim  == 'max': return max_value




def remove_prefix_from_keys(dictionary):
    new_dict = {}
    for key, value in dictionary.items():
        new_key = key.split('_')[-2] + "_" + key.split('_')[-1] if '_' in key else key  # Split at the first underscore and take the second part
        new_dict[new_key] = value
    return new_dict

# # Example usage:
# original_dict = {
#     'abc_key1': 10,
#     'def_key2': 20,
#     'ghi_key3': 30
# }

# new_dict = remove_prefix_from_keys(original_dict)
# print(new_dict)

    
#%%





